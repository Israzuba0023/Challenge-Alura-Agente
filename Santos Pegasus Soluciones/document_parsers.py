"""
Extratores de documentos multi-formato para o Agente Inteligente Corp - Santos Pegasus Soluciones.
Suporta os 8 formatos exigidos no Desafio Alura Agentes:
PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), Markdown (.md), CSV (.csv), JSON (.json) e HTML (.html).
"""

import os
import re
import csv
import json
from typing import List, Dict, Any, Optional

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None


class DocumentChunk:
    def __init__(
        self,
        content: str,
        document_name: str,
        category: str,
        file_format: str,
        page_number: Optional[int] = None,
        section: Optional[str] = None,
        chunk_id: str = ""
    ):
        self.content = content.strip()
        self.document_name = document_name
        self.category = category
        self.file_format = file_format
        self.page_number = page_number
        self.section = section
        self.chunk_id = chunk_id

    def to_dict(self) -> Dict[str, Any]:
        return {
            "chunk_id": self.chunk_id,
            "content": self.content,
            "document_name": self.document_name,
            "category": self.category,
            "file_format": self.file_format,
            "page_number": self.page_number,
            "section": self.section
        }


class MultiFormatDocumentParser:
    """Parser unificado para ler e chunkar múltiplos formatos de documento."""

    def detect_category(self, filename: str, content_sample: str = "") -> str:
        fn_lower = filename.lower()
        content_lower = content_sample.lower()

        if "back-end" in fn_lower or "backend" in fn_lower or "java" in content_lower or "microsserviço" in content_lower:
            return "Engenharia Back-end"
        elif "front-end" in fn_lower or "frontend" in fn_lower or "react" in content_lower or "css" in content_lower:
            return "Engenharia Front-end"
        elif "onboarding" in fn_lower or "desenvolvedores" in fn_lower or "boas-vindas" in content_lower:
            return "Onboarding & Pessoas"
        elif "arquitetura" in fn_lower or "domínios" in fn_lower or "dominio" in fn_lower:
            return "Arquitetura & Sistemas"
        elif "resiliência" in fn_lower or "incidentes" in fn_lower or "sop" in fn_lower or "p0" in content_lower:
            return "Resiliência & Operações"
        elif "benefícios" in fn_lower or "rh" in fn_lower:
            return "Recursos Humanos"
        elif "salários" in fn_lower or "cargos" in fn_lower or "faixa" in content_lower:
            return "Recursos Humanos & Financeiro"
        elif "reembolso" in fn_lower or "financeiro" in fn_lower:
            return "Financeiro e Contábil"
        elif "lgpd" in fn_lower or "compliance" in fn_lower or "termos" in fn_lower:
            return "Legal & Compliance"
        else:
            return "Geral Corporativo"

    def parse_file(self, file_path: str) -> List[DocumentChunk]:
        filename = os.path.basename(file_path)
        ext = os.path.splitext(filename)[1].lower()

        if ext == ".pdf":
            return self._parse_pdf(file_path, filename)
        elif ext == ".csv":
            return self._parse_csv(file_path, filename)
        elif ext == ".json":
            return self._parse_json(file_path, filename)
        elif ext in [".md", ".markdown"]:
            return self._parse_markdown(file_path, filename)
        elif ext in [".html", ".htm"]:
            return self._parse_html(file_path, filename)
        elif ext in [".txt"]:
            return self._parse_text(file_path, filename)
        elif ext in [".docx", ".doc"]:
            return self._parse_docx(file_path, filename)
        elif ext in [".xlsx", ".xls"]:
            return self._parse_xlsx(file_path, filename)
        elif ext in [".pptx", ".ppt"]:
            return self._parse_pptx(file_path, filename)
        else:
            return self._parse_text(file_path, filename)

    def _parse_pdf(self, file_path: str, filename: str) -> List[DocumentChunk]:
        chunks = []
        category = self.detect_category(filename)

        if pypdf is not None:
            try:
                reader = pypdf.PdfReader(file_path)
                for page_idx, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    if not page_text.strip():
                        continue
                    
                    page_chunks = self._chunk_text(
                        text=page_text,
                        doc_name=filename,
                        category=category,
                        fmt="PDF",
                        page_num=page_idx + 1
                    )
                    chunks.extend(page_chunks)
                return chunks
            except Exception as e:
                print(f"[Warning] Error parsing PDF with pypdf {filename}: {e}")

        # Fallback se pypdf falhar ou não extrair
        with open(file_path, "rb") as f:
            raw = f.read().decode("latin1", errors="ignore")
            clean_text = re.sub(r"[^\w\s\.\,\-\:\/\(\)\%\$\@]", " ", raw)
            clean_text = re.sub(r"\s+", " ", clean_text)
            if len(clean_text) > 100:
                chunks.extend(self._chunk_text(clean_text, filename, category, "PDF", page_num=1))

        return chunks

    def _parse_csv(self, file_path: str, filename: str) -> List[DocumentChunk]:
        chunks = []
        category = self.detect_category(filename)
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            rows = list(reader)

        if not rows:
            return chunks

        header = rows[0]
        rows_per_chunk = 5
        for i in range(1, len(rows), rows_per_chunk):
            batch = rows[i:i + rows_per_chunk]
            block_lines = [f"Colunas: {', '.join(header)}"]
            for r in batch:
                row_str = ", ".join([f"{header[j] if j < len(header) else 'Col'}: {val}" for j, val in enumerate(r)])
                block_lines.append(f"- {row_str}")
            
            content = "\n".join(block_lines)
            chunk = DocumentChunk(
                content=content,
                document_name=filename,
                category=category,
                file_format="CSV",
                section=f"Linhas {i} a {min(i + rows_per_chunk - 1, len(rows)-1)}"
            )
            chunks.append(chunk)

        return chunks

    def _parse_json(self, file_path: str, filename: str) -> List[DocumentChunk]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            try:
                data = json.load(f)
                content = json.dumps(data, indent=2, ensure_ascii=False)
            except Exception:
                f.seek(0)
                content = f.read()

        category = self.detect_category(filename, content)
        return self._chunk_text(content, filename, category, "JSON")

    def _parse_markdown(self, file_path: str, filename: str) -> List[DocumentChunk]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        category = self.detect_category(filename, text)
        return self._chunk_text(text, filename, category, "Markdown")

    def _parse_html(self, file_path: str, filename: str) -> List[DocumentChunk]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            raw_html = f.read()

        if BeautifulSoup is not None:
            soup = BeautifulSoup(raw_html, "html.parser")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.extract()
            text = soup.get_text(separator="\n")
        else:
            text = re.sub(r"<[^>]+>", " ", raw_html)

        text = re.sub(r"\n\s*\n", "\n\n", text)
        category = self.detect_category(filename, text)
        return self._chunk_text(text, filename, category, "HTML")

    def _parse_docx(self, file_path: str, filename: str) -> List[DocumentChunk]:
        try:
            import docx
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
        except Exception:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        category = self.detect_category(filename, text)
        return self._chunk_text(text, filename, category, "Word")

    def _parse_xlsx(self, file_path: str, filename: str) -> List[DocumentChunk]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            lines = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                lines.append(f"--- Planilha: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(val) for val in row if val is not None]
                    if row_vals:
                        lines.append(" | ".join(row_vals))
            text = "\n".join(lines)
        except Exception:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        category = self.detect_category(filename, text)
        return self._chunk_text(text, filename, category, "Excel")

    def _parse_pptx(self, file_path: str, filename: str) -> List[DocumentChunk]:
        chunks: List[DocumentChunk] = []
        category = self.detect_category(filename)
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                content = "\n".join(slide_text)
                if content.strip():
                    chunk = DocumentChunk(
                        content=content,
                        document_name=filename,
                        category=category,
                        file_format="PowerPoint",
                        page_number=idx + 1,
                        section=f"Slide {idx + 1}"
                    )
                    chunks.append(chunk)
            return chunks
        except Exception:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return self._chunk_text(text, filename, category, "PowerPoint")

    def _parse_text(self, file_path: str, filename: str) -> List[DocumentChunk]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        category = self.detect_category(filename, text)
        return self._chunk_text(text, filename, category, "Texto")

    def _chunk_text(
        self,
        text: str,
        doc_name: str,
        category: str,
        fmt: str,
        page_num: Optional[int] = None,
        max_chunk_size: int = 700
    ) -> List[DocumentChunk]:
        chunks = []
        paragraphs = text.split("\n\n")
        current_block = []
        current_len = 0

        for p in paragraphs:
            p_clean = p.strip()
            if not p_clean:
                continue

            if current_len + len(p_clean) > max_chunk_size and current_block:
                block_str = "\n\n".join(current_block)
                chunks.append(
                    DocumentChunk(
                        content=block_str,
                        document_name=doc_name,
                        category=category,
                        file_format=fmt,
                        page_number=page_num
                    )
                )
                current_block = [p_clean]
                current_len = len(p_clean)
            else:
                current_block.append(p_clean)
                current_len += len(p_clean)

        if current_block:
            block_str = "\n\n".join(current_block)
            chunks.append(
                DocumentChunk(
                    content=block_str,
                    document_name=doc_name,
                    category=category,
                    file_format=fmt,
                    page_number=page_num
                )
            )

        # Adicionar IDs únicos
        for idx, chk in enumerate(chunks):
            chk.chunk_id = f"{doc_name}_p{page_num or 1}_c{idx+1}"

        return chunks
